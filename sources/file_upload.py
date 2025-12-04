"""
File Upload Source Adapter

Handles file uploads for PDF, DOCX, TXT, and other document formats.
Extracts text content and metadata from uploaded files.

Supported formats:
    - PDF (.pdf) - via pdfplumber
    - Word (.docx) - via python-docx
    - Text (.txt, .md) - direct read
    - PowerPoint (.pptx) - via python-pptx

Example:
    >>> from sources.file_upload import FileUploadAdapter
    >>> adapter = FileUploadAdapter(tenant_id="tenant-123")
    >>> documents = await adapter.fetch(file_path="/path/to/document.pdf")
"""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional
import mimetypes
import logging

from models.document import RawDocument, DocumentSource
from sources.base import BaseSourceAdapter, FetchError, ValidationError

logger = logging.getLogger(__name__)


class FileUploadAdapter(BaseSourceAdapter):
    """Adapter for processing uploaded document files.

    Handles extraction of text content and metadata from various
    document formats including PDF, DOCX, TXT, and more.

    Attributes:
        max_file_size: Maximum file size in bytes (default: 50MB)
        extract_metadata: Whether to extract document metadata

    Example:
        >>> adapter = FileUploadAdapter(
        ...     tenant_id="tenant-123",
        ...     max_file_size=50 * 1024 * 1024
        ... )
        >>> docs = await adapter.fetch(file_path="report.pdf")
    """

    SUPPORTED_FORMATS = {
        ".pdf": "application/pdf",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }

    def __init__(
        self,
        tenant_id: Optional[str] = None,
        max_file_size: int = 50 * 1024 * 1024,  # 50MB
        extract_metadata: bool = True
    ):
        """Initialize file upload adapter.

        Args:
            tenant_id: Multi-tenant identifier
            max_file_size: Maximum file size in bytes
            extract_metadata: Whether to extract document metadata

        Example:
            >>> adapter = FileUploadAdapter(
            ...     tenant_id="tenant-123",
            ...     max_file_size=100 * 1024 * 1024  # 100MB
            ... )
        """
        super().__init__(
            source_type=DocumentSource.FILE_UPLOAD,
            tenant_id=tenant_id
        )
        self.max_file_size = max_file_size
        self.extract_metadata = extract_metadata

    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats.

        Returns:
            List of supported file extensions

        Example:
            >>> formats = adapter.get_supported_formats()
            >>> print(formats)
            ['.pdf', '.txt', '.md', '.docx', '.pptx']
        """
        return list(self.SUPPORTED_FORMATS.keys())

    async def validate_input(self, file_path: str, **kwargs) -> bool:
        """Validate file path and format.

        Args:
            file_path: Path to file
            **kwargs: Additional arguments (ignored)

        Returns:
            True if validation passes

        Raises:
            ValidationError: If file doesn't exist, too large, or unsupported format

        Example:
            >>> await adapter.validate_input("/path/to/doc.pdf")
            True
        """
        path = Path(file_path)

        # Check file exists
        if not path.exists():
            raise ValidationError(
                f"File does not exist: {file_path}",
                source=self.source_type.value,
                file_path=file_path
            )

        # Check file size
        file_size = path.stat().st_size
        if file_size > self.max_file_size:
            raise ValidationError(
                f"File too large: {file_size} bytes (max: {self.max_file_size})",
                source=self.source_type.value,
                file_path=file_path,
                file_size=file_size,
                max_size=self.max_file_size
            )

        # Check file format
        file_ext = path.suffix.lower()
        if file_ext not in self.SUPPORTED_FORMATS:
            raise ValidationError(
                f"Unsupported file format: {file_ext}",
                source=self.source_type.value,
                file_path=file_path,
                file_extension=file_ext,
                supported_formats=list(self.SUPPORTED_FORMATS.keys())
            )

        return True

    async def fetch(self, file_path: str, **kwargs) -> List[RawDocument]:
        """Fetch document from file path.

        Args:
            file_path: Path to document file
            **kwargs: Additional metadata to include

        Returns:
            List containing single RawDocument

        Raises:
            ValidationError: If file validation fails
            FetchError: If text extraction fails

        Example:
            >>> docs = await adapter.fetch(
            ...     file_path="/path/to/report.pdf",
            ...     author="John Doe"
            ... )
            >>> print(f"Extracted {len(docs[0].content)} characters")
        """
        # Validate input
        await self.validate_input(file_path)

        path = Path(file_path)
        file_ext = path.suffix.lower()

        self.logger.info(
            f"Fetching document from file: {file_path}",
            extra={
                "file_path": file_path,
                "file_ext": file_ext,
                "tenant_id": self.tenant_id
            }
        )

        try:
            # Extract content based on file type
            if file_ext == ".pdf":
                content, metadata = await self._extract_pdf(path)
            elif file_ext == ".docx":
                content, metadata = await self._extract_docx(path)
            elif file_ext == ".pptx":
                content, metadata = await self._extract_pptx(path)
            elif file_ext in [".txt", ".md"]:
                content, metadata = await self._extract_text(path)
            else:
                raise FetchError(
                    f"No handler for file type: {file_ext}",
                    source=self.source_type.value,
                    file_path=file_path
                )

            # Add file metadata
            file_metadata = {
                "filename": path.name,
                "file_ext": file_ext,
                "file_size": path.stat().st_size,
                "mime_type": self.SUPPORTED_FORMATS[file_ext],
                **metadata,
                **kwargs  # Include any additional metadata passed in
            }

            # Create RawDocument
            document = self._create_raw_document(
                content=content,
                url=f"file://{path.absolute()}",
                metadata=file_metadata
            )

            self.logger.info(
                f"Successfully extracted {len(content)} characters from {path.name}",
                extra={
                    "file_path": file_path,
                    "content_length": len(content),
                    "document_id": document.id,
                    "tenant_id": self.tenant_id
                }
            )

            return [document]

        except Exception as e:
            error_msg = f"Failed to extract content from {file_path}: {str(e)}"
            self.logger.error(
                error_msg,
                extra={
                    "file_path": file_path,
                    "error": str(e),
                    "tenant_id": self.tenant_id
                },
                exc_info=True
            )
            raise FetchError(
                error_msg,
                source=self.source_type.value,
                file_path=file_path,
                error=str(e)
            )

    async def _extract_pdf(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """Extract text and metadata from PDF file.

        Args:
            path: Path to PDF file

        Returns:
            Tuple of (content, metadata)

        Example:
            >>> content, metadata = await adapter._extract_pdf(Path("doc.pdf"))
            >>> print(f"Pages: {metadata['page_count']}")
        """
        try:
            import pdfplumber
        except ImportError:
            raise FetchError(
                "pdfplumber not installed. Install with: pip install pdfplumber",
                source=self.source_type.value
            )

        content_parts = []
        metadata: Dict[str, Any] = {}

        with pdfplumber.open(path) as pdf:
            # Extract metadata
            if self.extract_metadata and pdf.metadata:
                metadata.update({
                    "pdf_title": pdf.metadata.get("Title"),
                    "pdf_author": pdf.metadata.get("Author"),
                    "pdf_subject": pdf.metadata.get("Subject"),
                    "pdf_creator": pdf.metadata.get("Creator"),
                })

            metadata["page_count"] = len(pdf.pages)

            # Extract text from each page
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    content_parts.append(page_text)

                self.logger.debug(
                    f"Extracted {len(page_text) if page_text else 0} chars from page {i+1}",
                    extra={"page": i+1, "file": path.name}
                )

        content = "\n\n".join(content_parts)
        return content, metadata

    async def _extract_docx(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """Extract text and metadata from DOCX file.

        Args:
            path: Path to DOCX file

        Returns:
            Tuple of (content, metadata)

        Example:
            >>> content, metadata = await adapter._extract_docx(Path("doc.docx"))
            >>> print(f"Paragraphs: {metadata['paragraph_count']}")
        """
        try:
            from docx import Document
        except ImportError:
            raise FetchError(
                "python-docx not installed. Install with: pip install python-docx",
                source=self.source_type.value
            )

        doc = Document(path)
        content_parts = []
        metadata: Dict[str, Any] = {}

        # Extract metadata
        if self.extract_metadata and doc.core_properties:
            metadata.update({
                "docx_title": doc.core_properties.title,
                "docx_author": doc.core_properties.author,
                "docx_subject": doc.core_properties.subject,
                "docx_created": doc.core_properties.created.isoformat() if doc.core_properties.created else None,
            })

        # Extract text from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                content_parts.append(para.text)

        metadata["paragraph_count"] = len(content_parts)

        content = "\n\n".join(content_parts)
        return content, metadata

    async def _extract_pptx(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """Extract text and metadata from PPTX file.

        Args:
            path: Path to PPTX file

        Returns:
            Tuple of (content, metadata)

        Example:
            >>> content, metadata = await adapter._extract_pptx(Path("slides.pptx"))
            >>> print(f"Slides: {metadata['slide_count']}")
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise FetchError(
                "python-pptx not installed. Install with: pip install python-pptx",
                source=self.source_type.value
            )

        prs = Presentation(path)
        content_parts = []
        metadata: Dict[str, Any] = {}

        # Extract metadata
        if self.extract_metadata and prs.core_properties:
            metadata.update({
                "pptx_title": prs.core_properties.title,
                "pptx_author": prs.core_properties.author,
                "pptx_subject": prs.core_properties.subject,
            })

        metadata["slide_count"] = len(prs.slides)

        # Extract text from slides
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

            if slide_text:
                content_parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_text))

        content = "\n\n".join(content_parts)
        return content, metadata

    async def _extract_text(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """Extract text from plain text file.

        Args:
            path: Path to text file

        Returns:
            Tuple of (content, metadata)

        Example:
            >>> content, metadata = await adapter._extract_text(Path("doc.txt"))
            >>> print(f"Lines: {metadata['line_count']}")
        """
        # Try different encodings
        encodings = ["utf-8", "latin-1", "cp1252"]
        content = None
        encoding_used = None

        for encoding in encodings:
            try:
                with open(path, "r", encoding=encoding) as f:
                    content = f.read()
                encoding_used = encoding
                break
            except UnicodeDecodeError:
                continue

        if content is None:
            raise FetchError(
                f"Failed to decode file with encodings: {encodings}",
                source=self.source_type.value,
                file_path=str(path)
            )

        metadata = {
            "encoding": encoding_used,
            "line_count": content.count("\n") + 1,
        }

        return content, metadata


# Example usage
if __name__ == "__main__":
    import asyncio

    async def test_file_upload():
        """Test file upload adapter."""
        adapter = FileUploadAdapter(tenant_id="tenant-test")

        print(f"Supported formats: {adapter.get_supported_formats()}")

        # Create a test text file
        test_file = Path("test_document.txt")
        test_file.write_text("This is a test document.\n\nIt has multiple paragraphs.")

        try:
            # Fetch document
            docs = await adapter.fetch(
                file_path=str(test_file),
                author="Test Author"
            )

            print(f"\nFetched {len(docs)} document(s)")
            doc = docs[0]
            print(f"Document ID: {doc.id}")
            print(f"Content length: {len(doc.content)} chars")
            print(f"Metadata: {doc.metadata}")
            print(f"\nContent preview:\n{doc.content[:100]}...")

        finally:
            # Cleanup
            if test_file.exists():
                test_file.unlink()

    asyncio.run(test_file_upload())
